"""
VanGATE Website Analysis Report - Simplified Version
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# Create presentation
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Colors
DARK_BLUE = RGBColor(0x1A, 0x5F, 0xA5)
LIGHT_BLUE = RGBColor(0x37, 0x8A, 0xDD)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GRAY = RGBColor(0x5F, 0x5E, 0x5A)
LIGHT_GRAY = RGBColor(0xF1, 0xEF, 0xE8)
GREEN = RGBColor(0x63, 0x99, 0x22)
ORANGE = RGBColor(0xD8, 0x5A, 0x30)

def add_title_slide(prs, title, subtitle=""):
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    # Background
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = DARK_BLUE
    shape.line.fill.background()
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.8), Inches(12.333), Inches(1.5))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER
    
    # Subtitle
    if subtitle:
        sub_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.5), Inches(12.333), Inches(1))
        tf = sub_box.text_frame
        p = tf.paragraphs[0]
        p.text = subtitle
        p.font.size = Pt(22)
        p.font.color.rgb = RGBColor(0xB5, 0xD4, 0xF4)
        p.alignment = PP_ALIGN.CENTER
    
    return slide

def add_content_slide(prs, title, bullets):
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    # Header
    header = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.1))
    header.fill.solid()
    header.fill.fore_color.rgb = DARK_BLUE
    header.line.fill.background()
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.25), Inches(12.333), Inches(0.7))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = WHITE
    
    # Content
    content_box = slide.shapes.add_textbox(Inches(0.7), Inches(1.5), Inches(11.9), Inches(5.5))
    tf = content_box.text_frame
    tf.word_wrap = True
    
    for i, bullet in enumerate(bullets):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        
        if isinstance(bullet, dict):
            p.text = "• " + bullet.get('text', '')
            p.font.size = Pt(bullet.get('size', 20))
            p.font.color.rgb = RGBColor(0x44, 0x44, 0x41)
            if bullet.get('bold'):
                p.font.bold = True
        else:
            p.text = "• " + bullet
            p.font.size = Pt(20)
            p.font.color.rgb = RGBColor(0x44, 0x44, 0x41)
        p.space_after = Pt(14)
    
    return slide

def add_table_slide(prs, title, headers, rows, col_widths=None):
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    # Header
    header = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.1))
    header.fill.solid()
    header.fill.fore_color.rgb = DARK_BLUE
    header.line.fill.background()
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.25), Inches(12.333), Inches(0.7))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = WHITE
    
    # Table
    cols = len(headers)
    rows_count = len(rows) + 1
    table = slide.shapes.add_table(rows_count, cols, Inches(0.3), Inches(1.4), Inches(12.7), Inches(5.5)).table
    
    # Set column widths if provided
    if col_widths:
        for i, w in enumerate(col_widths):
            table.columns[i].width = Inches(w)
    
    # Header row
    for i, h in enumerate(headers):
        cell = table.cell(0, i)
        cell.text = h
        cell.fill.solid()
        cell.fill.fore_color.rgb = DARK_BLUE
        p = cell.text_frame.paragraphs[0]
        p.font.color.rgb = WHITE
        p.font.bold = True
        p.font.size = Pt(14)
        p.alignment = PP_ALIGN.CENTER
    
    # Data rows
    for row_idx, row in enumerate(rows):
        for col_idx, val in enumerate(row):
            cell = table.cell(row_idx + 1, col_idx)
            cell.text = str(val)
            p = cell.text_frame.paragraphs[0]
            p.font.size = Pt(12)
            p.font.color.rgb = RGBColor(0x44, 0x44, 0x41)
            p.alignment = PP_ALIGN.LEFT if col_idx == 0 else PP_ALIGN.CENTER
    
    return slide

def add_section_slide(prs, section_num, section_title):
    """Section divider slide"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    # Background
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = LIGHT_GRAY
    shape.line.fill.background()
    
    # Section number
    num_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(12.333), Inches(1))
    tf = num_box.text_frame
    p = tf.paragraphs[0]
    p.text = f"Part {section_num}"
    p.font.size = Pt(24)
    p.font.color.rgb = LIGHT_BLUE
    p.alignment = PP_ALIGN.CENTER
    
    # Section title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(3.2), Inches(12.333), Inches(1.2))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = section_title
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = DARK_BLUE
    p.alignment = PP_ALIGN.CENTER
    
    return slide

# ==================== SLIDES ====================

# Slide 1: Cover
add_title_slide(prs, "VanGATE Website Analysis", "Website Functionality Overview | 2026-06-11")

# Slide 2: Site Overview
add_content_slide(prs, "Website Overview", [
    {"text": "Website URL: https://vangate.ca/", "size": 20},
    {"text": "Company: VanGATE Manufacturing & Automation", "size": 20},
    {"text": "Location: Richmond, BC, Canada", "size": 20},
    {"text": "Industry: Automatic Door & Gate Systems", "size": 20},
    {"text": "Current Platform: WordPress 6.9.4", "size": 20},
    {"text": "Business Focus: Residential & Commercial Solutions", "size": 20},
    {"text": " ", "size": 12},
    {"text": "Website Purpose:", "size": 20, "bold": True},
    "Introduce company and products/services",
    "Showcase products, projects, and galleries",
    "Publish news and industry articles",
    "Display customer testimonials",
    "Receive customer inquiries via form and phone",
    "Support online marketing via SEO and analytics"
])

# Slide 3: User Roles
add_table_slide(prs, "User Roles",
    ["Role", "Current Functions", "Status"],
    [
        ["Website Visitor", "Browse content, view products/services, gallery, news; submit inquiry, call", "Confirmed"],
        ["Content Editor", "Manage pages, articles, images, menus, SEO info", "Assumed"],
        ["Administrator", "Manage users, themes, plugins, forms, settings", "Assumed"]
    ],
    col_widths=[2.5, 7, 2.5]
)

# ==================== Part 1: 功能说明总览 ====================
add_section_slide(prs, "1", "Functionality Overview")

# Slide: Overview Summary
add_table_slide(prs, "Functionality Summary",
    ["Module", "Main Features", "Status"],
    [
        ["Navigation", "Desktop top nav, mobile menu, logo, footer links", "Confirmed"],
        ["Homepage", "Hero, products overview, gallery, testimonials, news, CTA", "Confirmed"],
        ["Company", "Company introduction page", "Confirmed"],
        ["Products/Services", "Category listing, product detail pages", "Confirmed"],
        ["Gallery", "Project photos with lightbox", "Confirmed"],
        ["News/Articles", "Article listing, detail pages, pagination", "Confirmed"],
        ["Contact", "Contact form, phone, address, department info", "Confirmed"],
        ["CMS Backend", "WordPress admin, content management", "Assumed"]
    ],
    col_widths=[2.5, 7.5, 2.2]
)

# ==================== Part 2: 功能说明列表 ====================
add_section_slide(prs, "2", "Detailed Functionality List")

# Slide: Navigation
add_table_slide(prs, "Navigation & Public Areas (WEB Series)",
    ["ID", "Function", "Status"],
    [
        ["WEB-001", "Desktop top navigation bar", "Confirmed"],
        ["WEB-002", "Mobile menu", "Confirmed"],
        ["WEB-003", "VanGATE logo with home link", "Confirmed"],
        ["WEB-004", "Navigation entries: Products, Services, Company, News, Gallery, Contact", "Confirmed"],
        ["WEB-005", "Clickable phone number", "Confirmed"],
        ["WEB-006", "Footer with services, products, company, address, phone", "Confirmed"],
        ["WEB-007", "Back to top button", "Confirmed"],
        ["WEB-008", "Responsive design (desktop, tablet, mobile)", "Confirmed"]
    ],
    col_widths=[1.2, 8.5, 2.3]
)

# Slide: Homepage
add_table_slide(prs, "Homepage (HOM Series)",
    ["ID", "Function", "Status"],
    [
        ["HOM-001", "Introduction to residential & commercial solutions", "Confirmed"],
        ["HOM-002", "Primary CTA: contact and call buttons", "Confirmed"],
        ["HOM-003", "Main products and services overview", "Confirmed"],
        ["HOM-004", "Project/gallery image showcase", "Confirmed"],
        ["HOM-005", "Gallery images can be enlarged", "Confirmed"],
        ["HOM-006", "Link to full gallery", "Confirmed"],
        ["HOM-007", "Customer testimonials section", "Confirmed"],
        ["HOM-008", "Recent news/articles", "Confirmed"],
        ["HOM-009", "Link to full news list", "Confirmed"],
        ["HOM-010", "Effects: carousel, lightbox, scroll animations", "Confirmed"]
    ],
    col_widths=[1.2, 8.5, 2.3]
)

# Slide: Products & Services
add_table_slide(prs, "Products & Services Catalog (CAT Series)",
    ["ID", "Function", "Status"],
    [
        ["CAT-001", "Products and services organized by category", "Confirmed"],
        ["CAT-002", "Category listing pages", "Confirmed"],
        ["CAT-003", "Individual product/service detail pages", "Confirmed"],
        ["CAT-004", "Content includes: name, description, images", "Confirmed"],
        ["CAT-005", "Each product/service has unique URL", "Confirmed"],
        ["CAT-006", "Product pages include contact CTA", "Confirmed"],
        ["CAT-007", "Catalog covers both residential and commercial", "Confirmed"],
        ["CAT-008", "Backend category and content maintenance", "Assumed"]
    ],
    col_widths=[1.2, 8.5, 2.3]
)

# Slide: Gallery & News
add_table_slide(prs, "Gallery & News (GAL/NEW Series)",
    ["ID", "Function", "Status"],
    [
        ["GAL-001", "Project/product gallery", "Confirmed"],
        ["GAL-002", "Gallery displayed as image collection", "Confirmed"],
        ["GAL-003", "Gallery responsive on all screen sizes", "Confirmed"],
        ["GAL-004", "Images can be enlarged, browse prev/next", "Confirmed"],
        ["GAL-005", "Backend image upload and reuse", "Assumed"],
        ["NEW-001", "News or article listing", "Confirmed"],
        ["NEW-002", "List shows: title, date, summary, image", "Confirmed"],
        ["NEW-003", "Each article has dedicated detail page", "Confirmed"],
        ["NEW-004", "News list supports pagination", "Confirmed"],
        ["NEW-005", "Backend article CRUD operations", "Assumed"]
    ],
    col_widths=[1.2, 8.5, 2.3]
)

# Slide: Contact
add_table_slide(prs, "Contact & Inquiries (CON Series)",
    ["ID", "Function", "Status"],
    [
        ["CON-001", "Contact page", "Confirmed"],
        ["CON-002", "Company address and phone display", "Confirmed"],
        ["CON-003", "Clickable phone number on mobile", "Confirmed"],
        ["CON-004", "Online inquiry form", "Confirmed"],
        ["CON-005", "Form validates required fields", "Confirmed"],
        ["CON-006", "Form submission triggers analytics event", "Confirmed"],
        ["CON-007", "Form submission sends email to recipients", "Assumed"],
        ["CON-008", "Form fields, recipients, spam protection details", "To Confirm"]
    ],
    col_widths=[1.2, 8.5, 2.3]
)

# Slide: SEO & Analytics
add_table_slide(prs, "SEO & Analytics (SEO Series)",
    ["ID", "Function", "Status"],
    [
        ["SEO-001", "Page supports independent SEO title and description", "Confirmed"],
        ["SEO-002", "Page outputs canonical URLs", "Confirmed"],
        ["SEO-003", "Page outputs Open Graph and Twitter card info", "Confirmed"],
        ["SEO-004", "Page outputs structured data (JSON-LD)", "Confirmed"],
        ["SEO-005", "Website provides XML sitemap", "Assumed"],
        ["SEO-006", "Google Analytics integration", "Confirmed"],
        ["SEO-007", "Google Tag Manager integration", "Confirmed"],
        ["SEO-008", "Contact form as analytics conversion event", "Confirmed"]
    ],
    col_widths=[1.2, 8.5, 2.3]
)

# Slide: Technical Stack
add_table_slide(prs, "Technical Environment",
    ["Component", "Current Implementation", "Status"],
    [
        ["Web Server", "LiteSpeed", "Confirmed"],
        ["Server Language", "PHP 7.4.33", "Confirmed"],
        ["CMS", "WordPress 6.9.4", "Confirmed"],
        ["Theme", "Custom 'vangate' theme", "Confirmed"],
        ["Page Editor", "Elementor", "Confirmed"],
        ["Frontend", "Bootstrap 4.5.3, jQuery, Custom CSS", "Confirmed"],
        ["Contact Form", "Contact Form 7", "Confirmed"],
        ["SEO Plugin", "Yoast SEO", "Confirmed"],
        ["Security", "Wordfence", "Confirmed"],
        ["Components", "Swiper, Slick, PhotoSwipe, Lity, WOW.js, Font Awesome", "Confirmed"]
    ],
    col_widths=[2.5, 6, 2.5]
)

# Slide: Not Observed
add_content_slide(prs, "Functions NOT Observed", [
    "Customer registration and login",
    "Dealer or contractor portal",
    "Online quote tool",
    "Shopping cart and online payment",
    "Inventory and warehouse management",
    "Installation project progress tracking",
    "After-sales ticket management",
    "Online appointment booking",
    "CRM synchronization",
    "Multi-language switching",
    "Employee internal system"
])

# Slide: Confirmation Items
add_content_slide(prs, "Items to Confirm", [
    "Which pages and homepage modules can clients edit?",
    "Products/services: custom post type or theme template?",
    "Does gallery and testimonials have independent admin entry?",
    "What fields does contact form include, who receives emails?",
    "Are contact form records saved to database?",
    "What admin users and roles exist?",
    "Are there unpublished pages still in use?",
    "Are article categories, tags, authors, comments used?",
    "Who manages analytics and search console accounts?",
    "Is there regular backup and security update schedule?"
])

# ==================== Part 3: 新网站设计风格 ====================
add_section_slide(prs, "3", "New Website Design Styles")

# Slide: Design Style Overview
add_table_slide(prs, "Three Design Styles Overview",
    ["Style", "Positioning", "Target Audience", "Key Features"],
    [
        ["Modern Minimal\n现代极简", "Clean & Professional\n专业清新", "High-end homeowners\nRepair chains\nB2B clients", "Light theme, ample whitespace\nSoft shadows, refined details\nTrust and professionalism"],
        ["Industrial Professional\n工业专业", "Heavy-Duty & Bold\n稳重功能", "Commercial clients\nRepair technicians\nFleet managers", "Dark theme, industrial orange\nBold typography\nData-focused layout"],
        ["Vibrant Contemporary\n活力现代", "Bold & Dynamic\n年轻动感", "Young car owners\nAuto enthusiasts\nE-commerce retail", "Gradient colors, micro-animations\nCard-based layout\nEnergetic and modern"]
    ],
    col_widths=[2.5, 2.8, 3, 4.2]
)

# Slide: Style 1 - Modern Minimal
add_table_slide(prs, "Style 1: Modern Minimal (现代极简)",
    ["Aspect", "Specification"],
    [
        ["Color Palette", "Primary: #1E3A5F (Deep Sea Blue)\nAccent: #0EA5E9 (Sky Blue)\nBackground: #FFFFFF\nText: #1E293B"],
        ["Typography", "Headings: Plus Jakarta Sans (600, 700)\nBody: Inter (400, 500)\nStrong contrast between heading and body weights"],
        ["Design Elements", "Generous whitespace\nSoft shadows (0 4px 12px)\nRounded corners (10-16px)\nMinimal borders"],
        ["Mobile Design", "Single column layout\n64px fixed header\nFull-width cards\nBottom-anchored CTAs"],
        ["PC Design", "Max-width: 1280px centered\n6-column category grid\n4-column product grid\nHover animations on cards"]
    ],
    col_widths=[2.5, 10]
)

# Slide: Style 2 - Industrial Professional
add_table_slide(prs, "Style 2: Industrial Professional (工业专业)",
    ["Aspect", "Specification"],
    [
        ["Color Palette", "Primary: #F97316 (Industrial Orange)\nBackground: #18181B (Deep Black)\nSurface: #27272A (Carbon Gray)\nAccent: #FBBF24 (Amber Yellow)"],
        ["Typography", "Headings: Bebas Neue (700)\nBody: Source Sans 3 (400, 600)\nALL CAPS for key labels\nStrong industrial feel"],
        ["Design Elements", "Bold 2px borders\nSharp corners (2-8px radius)\nHigh contrast elements\nStats bar with prominent numbers"],
        ["Mobile Design", "Full-width dark sections\nOrange accent borders\nDiagonal background cuts\nUppercase navigation"],
        ["PC Design", "1200px max-width\nDiagonal skew decorations\n6-column category grid\n4-column product grid"]
    ],
    col_widths=[2.5, 10]
)

# Slide: Style 3 - Vibrant Contemporary
add_table_slide(prs, "Style 3: Vibrant Contemporary (活力现代)",
    ["Aspect", "Specification"],
    [
        ["Color Palette", "Primary: #2563EB (Electric Blue)\nSecondary: #DC2626 (Racing Red)\nAccent: #10B981 (Speed Green)\nGradient: 135deg Blue to Purple"],
        ["Typography", "Headings: Outfit (600, 700)\nBody: DM Sans (400, 500)\nGeometric modern feel\nRounded and approachable"],
        ["Design Elements", "Gradient backgrounds\nCard hover animations\nFloating action buttons\nFeature pills and badges"],
        ["Mobile Design", "Frosted glass nav bar\nFeature bar scroll\nGradient top accents\nHeart wishlist buttons"],
        ["PC Design", "1280px max-width\nHover lift effects\nGradient CTA sections\nAdd-to-cart buttons"]
    ],
    col_widths=[2.5, 10]
)

# Slide: Mobile vs PC Comparison
add_table_slide(prs, "Mobile vs PC Design Adaptation",
    ["Element", "Mobile (375px)", "Tablet (640px)", "Desktop (1024px+)"],
    [
        ["Navigation", "Hamburger menu\n64px header", "Compact horizontal\nReduced labels", "Full horizontal\nWith labels"],
        ["Grid Layout", "2 columns", "3 columns", "4-6 columns"],
        ["Product Cards", "Full-width stacked", "Grid with gaps", "Dense grid"],
        ["Hero Section", "Single column\nCentered text", "Left-aligned\nWith sidebar", "Full layout\nWith stats bar"],
        ["Touch Targets", "44px minimum", "44px minimum", "Hover states"],
        ["Footer", "Stacked links\nSingle column", "2-column grid", "3-column grid"]
    ],
    col_widths=[2.5, 3.2, 3.2, 3.6]
)

# Slide: Technical Specifications
add_table_slide(prs, "Design System Technical Specs",
    ["Specification", "Value"],
    [
        ["Mobile-First", "Base styles for 320px, progressive enhancement"],
        ["Responsive Breakpoints", "640px (tablet), 1024px (desktop), 1280px (large)"],
        ["Touch Targets", "Minimum 44x44px for all interactive elements"],
        ["Font Sizes", "Fluid typography with clamp() for headings"],
        ["Spacing System", "8px base unit (4, 8, 12, 16, 24, 32, 48, 64px)"],
        ["Color Contrast", "WCAG AA: 4.5:1 for normal text, 3:1 for large text"],
        ["Safe Areas", "viewport-fit=cover for notch handling"],
        ["Performance", "Lazy loading images, optimized font loading"]
    ],
    col_widths=[3, 9.5]
)

# Slide: Recommendation Summary
add_content_slide(prs, "Style Selection Recommendation", [
    {"text": "Based on VanGATE's B2B focus and North American market:", "size": 18, "bold": True},
    {"text": " ", "size": 10},
    {"text": "Modern Minimal (推荐)", "size": 22, "bold": True},
    "Best for: Professional image, trust-building, B2B clients",
    "Clean interface conveys reliability and quality",
    " ",
    {"text": "Industrial Professional", "size": 20, "bold": True},
    "Best for: Commercial/technical customers, fleet managers",
    "Emphasizes functionality and industrial-grade quality",
    " ",
    {"text": "Vibrant Contemporary", "size": 20, "bold": True},
    "Best for: E-commerce retail, younger demographic",
    "Modern and energetic, good for online sales emphasis"
])

# Slide: Summary
add_title_slide(prs, "Design Styles Summary", "Mobile-First | Responsive | WCAG AA Compliant | North American Aesthetic")

# Save
output_path = "D:/Vangate website/docs/VanGATE_Website_Analysis_Report_v2.pptx"
prs.save(output_path)
print(f"Presentation saved to: {output_path}")
